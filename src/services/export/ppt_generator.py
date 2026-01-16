import asyncio
import re
import json
from pathlib import Path
from datetime import datetime
from typing import Any, Dict, List, Optional, Tuple, Union

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.util import Pt
    from pptx.slide import Slide
    from pptx.shapes.autoshape import Shape
except ImportError:
    Presentation = None

from src.logging import get_logger
from src.services.config import get_ppt_config
from src.services.llm import complete as llm_complete


logger = get_logger("PPTGenerator")


class PPTGenerator:
    """
    Service for generating PowerPoint presentations from Markdown content.
    """

    def __init__(self, export_dir: Union[str, Path]):
        self.export_dir = Path(export_dir)
        self.export_dir.mkdir(parents=True, exist_ok=True)
        if Presentation is None:
            logger.warning("python-pptx not installed. PPT export will fail.")

    async def generate(
        self,
        markdown: str,
        title: Optional[str] = None,
        style_prompt: Optional[str] = None,
        style_model: Optional[str] = None,
        style_api_key: Optional[str] = None,
        style_base_url: Optional[str] = None,
        max_slides: int = 15,
    ) -> Dict[str, Any]:
        """
        Generate a PPTX file from markdown.

        Args:
            markdown: The markdown content.
            title: Optional title for the presentation.
            style_prompt: Optional style instructions for the LLM.
            style_model: Optional model override (uses PPT config if not provided).
            style_api_key: Optional API key override (uses PPT config if not provided).
            style_base_url: Optional base URL override (uses PPT config if not provided).
            max_slides: Maximum number of slides.

        Returns:
            Dict containing filename, relative_path, and download_url.
        """
        if Presentation is None:
            raise ImportError("python-pptx is required for PPT export.")

        # 1. Parse Content
        extracted_title, sections = self._split_markdown_into_sections(markdown)
        final_title = (title or "").strip() or extracted_title

        # 2. Prepare Filename (Fixing the sanitization bug here)
        safe_title = self._sanitize_filename(final_title)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{safe_title}_{timestamp}.pptx"
        output_path = self.export_dir / filename

        # 3. Generate Style/Structure Specification via LLM if requested
        spec = None
        if style_prompt:
            try:
                spec = await self._generate_ppt_spec(
                    markdown, style_prompt, max_slides,
                    model_override=style_model,
                    api_key_override=style_api_key,
                    base_url_override=style_base_url,
                )
            except Exception as e:
                logger.error(f"Failed to generate PPT spec via LLM: {e}")
                # Fallback to default rendering if LLM fails
                spec = None

        # 4. Create Presentation
        prs = Presentation()
        self._build_presentation(prs, final_title, sections, spec, max_slides)

        # 5. Save
        prs.save(str(output_path))
        logger.info(f"PPT generated: {output_path}")

        # Assuming the export_dir is under data/user/..., construct relative path
        # logic matches previous router implementation
        # project_root/data/user/research/exports -> relative: research/exports/filename
        # We'll try to determine relative path dynamically if possible, or assume standard structure
        try:
            # Try to find 'data' in path components
            parts = output_path.parts
            if "data" in parts and "user" in parts:
                user_idx = parts.index("user")
                # path after 'user'
                relative_path = "/".join(parts[user_idx + 1 :])
            else:
                # Fallback
                relative_path = f"research/exports/{filename}"
        except ValueError:
            relative_path = f"research/exports/{filename}"

        return {
            "filename": filename,
            "relative_path": relative_path,
            "download_url": f"/api/outputs/{relative_path}",
        }

    def _sanitize_filename(self, value: str) -> str:
        """
        Sanitize filename, allowing Unicode characters (Chinese, etc.).
        Only strip strictly invalid filesystem characters.
        """
        # Remove null bytes
        value = value.replace("\0", "")
        # Remove invalid chars for Windows/Unix: / \ : * ? " < > |
        # Also remove newlines just in case
        cleaned = re.sub(r'[\\/*?:"<>|\n\r]+', "", value)
        cleaned = cleaned.strip()
        return cleaned or "presentation"

    def _split_markdown_into_sections(
        self, markdown: str
    ) -> Tuple[str, List[Tuple[str, str]]]:
        """Split markdown into title and titled sections."""
        lines = markdown.splitlines()
        title = None
        sections = []

        current_heading = None
        current_lines = []

        for raw_line in lines:
            line = raw_line.strip()
            # Match headers # ... to ###### ...
            m = re.match(r"^(#{1,6})\s+(.*)$", line)
            if m:
                level = len(m.group(1))
                heading = m.group(2).strip()
                if level == 1 and title is None:
                    title = heading
                    continue

                # Treat H2 as section separators (flexible based on structure)
                # Adjusting logic to capture H2 as main slide sections
                if level <= 2:
                    if current_heading is not None:
                        sections.append(
                            (current_heading, "\n".join(current_lines).strip())
                        )
                    current_heading = heading
                    current_lines = []
                    continue

            if current_heading is not None:
                current_lines.append(raw_line)

        # Add last section
        if current_heading is not None:
            sections.append((current_heading, "\n".join(current_lines).strip()))

        if not title:
            # Try to use first section as title if H1 missing
            title = sections[0][0] if sections else "Research Report"

        return title, sections

    async def _generate_ppt_spec(
        self,
        markdown: str,
        style_prompt: str,
        max_slides: int,
        model_override: Optional[str] = None,
        api_key_override: Optional[str] = None,
        base_url_override: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Call LLM to generate structure and theme.
        
        Uses PPT-specific config with optional per-request overrides.
        """
        ppt_config = get_ppt_config()
        
        # Apply overrides (request params > ppt config)
        effective_model = model_override or ppt_config.model
        effective_api_key = api_key_override or ppt_config.api_key
        effective_base_url = base_url_override or ppt_config.base_url
        
        if not effective_model:
            logger.warning("No PPT model configured, skipping LLM spec generation")
            return {}
        
        system_prompt = (
            "You are an expert presentation designer. "
            "Return ONLY valid JSON with keys: title, theme, slides. "
            "theme must contain: background, accent, title_color, body_color, font. "
            "slides must be an array of objects: {title: string, bullets: string[]}. "
            "No markdown, no commentary."
        )
        user_prompt = (
            "Create a slide deck specification from the following markdown report.\n"
            f"Constraints:\n- Maximum slides: {max(1, max_slides)}\n"
            "- Keep bullets concise (<= 12 words each)\n"
            "- Ensure the style matches the following instruction:\n"
            f"{style_prompt}\n\n"
            "Report markdown:\n"
            f"{markdown}"
        )

        logger.info(f"Using PPT model: {effective_model}")
        raw = await llm_complete(
            prompt=user_prompt,
            system_prompt=system_prompt,
            model=effective_model,
            api_key=effective_api_key,
            base_url=effective_base_url,
            binding=ppt_config.binding,
            temperature=ppt_config.temperature,
            max_tokens=ppt_config.max_tokens,
            response_format={"type": "json_object"}
        )
        
        return self._extract_json(raw) or {}

    def _extract_json(self, text: str) -> Optional[Dict[str, Any]]:
        try:
            # Try direct parse
            return json.loads(text)
        except json.JSONDecodeError:
            pass
            
        # Try extracting code block
        m = re.search(r"```json\s*([\s\S]*?)\s*```", text, flags=re.IGNORECASE)
        if m:
            try:
                return json.loads(m.group(1).strip())
            except:
                pass
        
        # Try finding braces
        start = text.find("{")
        end = text.rfind("}")
        if start != -1 and end != -1 and end > start:
            try:
                return json.loads(text[start : end + 1])
            except:
                pass
        return None

    def _build_presentation(
        self,
        prs,
        title: str,
        sections: List[Tuple[str, str]],
        spec: Optional[Dict[str, Any]],
        max_slides: int,
    ):
        """Builds the PPT slides based on spec or fallback to sections."""
        
        # 1. Determine Theme
        theme = spec.get("theme", {}) if spec else {}
        theme_config = self._parse_theme(theme)
        
        # 2. Title Slide
        self._add_title_slide(prs, title, theme_config)

        # 3. Content Slides
        slide_budget = max_slides - 1  # Minus title slide
        bullet_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

        # Use spec slides if available
        slides_to_create = []
        if spec and isinstance(spec.get("slides"), list):
            for s in spec["slides"]:
                if isinstance(s, dict) and s.get("title") and isinstance(s.get("bullets"), list):
                    slides_to_create.append({
                        "title": s["title"],
                        "bullets": [str(b) for b in s["bullets"] if b]
                    })
        
        # If no spec slides, fallback to parsing sections
        if not slides_to_create:
            for sec_title, sec_body in sections:
                bullets = self._extract_bullets(sec_body)
                if not bullets:
                    bullets = [sec_body.strip()] if sec_body.strip() else []
                if not bullets:
                    continue
                
                # Chunk large sections
                chunk_size = 7
                chunks = [bullets[i : i + chunk_size] for i in range(0, len(bullets), chunk_size)]
                for idx, chunk in enumerate(chunks):
                    t = sec_title if idx == 0 else f"{sec_title} (cont.)"
                    slides_to_create.append({"title": t, "bullets": chunk})

        # Render slides
        for slide_data in slides_to_create:
            if slide_budget <= 0:
                break
            
            slide = prs.slides.add_slide(bullet_layout)
            self._apply_slide_style(slide, prs, theme_config)
            
            # Set Title
            if slide.shapes.title:
                slide.shapes.title.text = slide_data["title"]
                self._style_text(slide.shapes.title.text_frame, theme_config, is_title=True)
            
            # Set Body
            body_shape = self._find_body_placeholder(slide)
            if body_shape:
                 self._set_bullets(body_shape, slide_data["bullets"], theme_config)
            
            slide_budget -= 1

    def _parse_theme(self, theme: Dict[str, Any]) -> Dict[str, Any]:
        """Parse theme colors and font."""
        return {
            "background": self._parse_hex_color(str(theme.get("background", "#FFFFFF"))) or (255, 255, 255),
            "accent": self._parse_hex_color(str(theme.get("accent", "#4F46E5"))) or (79, 70, 229),
            "title_color": self._parse_hex_color(str(theme.get("title_color", "#111827"))) or (17, 24, 39),
            "body_color": self._parse_hex_color(str(theme.get("body_color", "#111827"))) or (17, 24, 39),
            "font": str(theme.get("font", "")).strip() or "Arial", # Default to safe font
        }

    def _parse_hex_color(self, value: str) -> Optional[Tuple[int, int, int]]:
        s = (value or "").strip().lstrip("#")
        if len(s) == 3:
            s = "".join([c * 2 for c in s])
        if len(s) != 6:
            return None
        try:
            return tuple(int(s[i:i+2], 16) for i in (0, 2, 4))
        except:
            return None

    def _add_title_slide(self, prs, title: str, theme: Dict[str, Any]):
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        
        # Background
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*theme["background"])
        
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._style_text(slide.shapes.title.text_frame, theme, is_title=True, font_size=44)
            
        # Clear subtitle if exists
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            if subtitle:
                subtitle.text = ""

    def _apply_slide_style(self, slide, prs, theme: Dict[str, Any]):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*theme["background"])
        
        # Add accent bar
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            0, 0, prs.slide_width, int(prs.slide_height * 0.05)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*theme["accent"])
        bar.line.fill.background()

    def _style_text(self, text_frame, theme: Dict[str, Any], is_title=False, font_size=None):
        if not text_frame.paragraphs:
            return
            
        # If font_size not provided, set defaults
        if font_size is None:
            font_size = 32 if is_title else 20
            
        color = theme["title_color"] if is_title else theme["body_color"]
        
        # Style existing paragraphs or first one
        for p in text_frame.paragraphs:
            p.font.name = theme["font"]
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor(*color)

    def _set_bullets(self, body_shape, bullets: List[str], theme: Dict[str, Any]):
        if not hasattr(body_shape, "text_frame"):
            return
        tf = body_shape.text_frame
        tf.clear()
        
        for i, item in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            
            p.font.name = theme["font"]
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(*theme["body_color"])

    def _find_body_placeholder(self, slide):
        for shape in slide.placeholders:
            if shape.is_placeholder and shape.placeholder_format.type == 2: # Body
                return shape
        # Fallback to index 1
        if len(slide.placeholders) > 1:
            return slide.placeholders[1]
        return None

    def _extract_bullets(self, text: str) -> List[str]:
        """Simple bullet extractor."""
        if not text:
            return []
        
        bullets = []
        for line in text.splitlines():
            line = line.strip()
            if not line: 
                continue
            # Match bullet markers
            if re.match(r"^[-*+]\s+", line) or re.match(r"^\d+\.\s+", line):
                clean = re.sub(r"^([-*+]|\d+\.)\s+", "", line)
                bullets.append(clean)
            else:
                 # If line is short enough, treat as bullet; otherwise maybe paragraph
                 if len(line.split()) < 30:
                     bullets.append(line)
        return bullets
