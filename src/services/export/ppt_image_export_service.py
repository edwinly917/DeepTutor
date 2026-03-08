from __future__ import annotations

from datetime import datetime
from pathlib import Path
import re
from typing import Any

from PIL import Image

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    Presentation = None
    Inches = None

from src.logging import get_logger

logger = get_logger("PPTImageExportService")


class PptImageExportService:
    def __init__(self, project_root: Path):
        self.project_root = project_root

    def export_pptx(
        self,
        *,
        project_id: str,
        title: str,
        image_paths: list[str],
        aspect_ratio: str = "16:9",
    ) -> dict[str, Any]:
        if Presentation is None or Inches is None:
            raise ImportError("python-pptx is required for PPT export.")

        export_dir = self._project_export_dir(project_id)
        filename = f"{self._sanitize_filename(title)}_{self._timestamp()}.pptx"
        output_path = export_dir / filename

        prs = Presentation()
        if aspect_ratio == "4:3":
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
        else:
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]
        slide_width_inches = prs.slide_width / 914400
        slide_height_inches = prs.slide_height / 914400

        for image_path in image_paths:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                image_path,
                0,
                0,
                width=Inches(slide_width_inches),
                height=Inches(slide_height_inches),
            )

        prs.save(str(output_path))
        return self._build_output_payload(output_path)

    def export_pdf(
        self,
        *,
        project_id: str,
        title: str,
        image_paths: list[str],
    ) -> dict[str, Any]:
        if not image_paths:
            raise ValueError("No image paths provided for PDF export")

        export_dir = self._project_export_dir(project_id)
        filename = f"{self._sanitize_filename(title)}_{self._timestamp()}.pdf"
        output_path = export_dir / filename

        images = []
        for image_path in image_paths:
            image = Image.open(image_path)
            if image.mode != "RGB":
                image = image.convert("RGB")
            images.append(image)

        first, rest = images[0], images[1:]
        first.save(str(output_path), save_all=True, append_images=rest)
        return self._build_output_payload(output_path)

    def _project_export_dir(self, project_id: str) -> Path:
        export_dir = (
            self.project_root / "data" / "user" / "ppt" / "projects" / project_id / "exports"
        )
        export_dir.mkdir(parents=True, exist_ok=True)
        return export_dir

    def _timestamp(self) -> str:
        return datetime.now().strftime("%Y%m%d_%H%M%S")

    def _sanitize_filename(self, value: str) -> str:
        cleaned = re.sub(r'[\\/*?:"<>|\n\r]+', "", (value or "").replace("\0", "")).strip()
        return cleaned or "presentation"

    def _build_output_payload(self, output_path: Path) -> dict[str, Any]:
        parts = output_path.parts
        user_idx = parts.index("user")
        relative_path = "/".join(parts[user_idx + 1 :])
        return {
            "filename": output_path.name,
            "relative_path": relative_path,
            "download_url": f"/api/outputs/{relative_path}",
        }
