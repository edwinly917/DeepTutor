from __future__ import annotations

import base64
from dataclasses import dataclass
from io import BytesIO
import json
from pathlib import Path
import re
import tempfile
from typing import Any

from PIL import Image, ImageDraw, ImageFont

from src.logging import get_logger
from src.services.config import get_ppt_analysis_vision_config
from src.services.llm import complete as llm_complete
from src.services.llm import get_token_limit_kwargs
from src.services.llm.client import LLMClient
from src.services.llm.config import LLMConfig, supports_response_format_json_object
from src.services.ppt.prompts import PptPromptManager

logger = get_logger("PptTemplateAnalyzer")

_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp"}
_PREVIEW_MAX_PAGES = 3
_PREVIEW_WIDTH = 1600
_DEFAULT_BG = (255, 255, 255)
_DEFAULT_SHAPE_FILL = (247, 249, 252)
_DEFAULT_SHAPE_LINE = (203, 213, 225)
_DEFAULT_TEXT = (15, 23, 42)


@dataclass
class TemplateAnalysisResult:
    reference_style_prompt: str = ""
    reference_layout_prompt: str = ""
    reference_content_prompt: str = ""
    extracted_text: str = ""

    def to_dict(self) -> dict[str, str]:
        return {
            "reference_style_prompt": self.reference_style_prompt,
            "reference_layout_prompt": self.reference_layout_prompt,
            "reference_content_prompt": self.reference_content_prompt,
            "extracted_text": self.extracted_text,
        }


class TemplateAnalyzer:
    def __init__(self, project_root: Path):
        self.project_root = project_root
        self.config = get_ppt_analysis_vision_config(project_root)
        self._vision_client = LLMClient(
            LLMConfig(
                model=self.config.model,
                api_key=self.config.api_key,
                base_url=self.config.base_url,
                binding=self.config.binding,
                temperature=self.config.temperature,
                max_tokens=self.config.max_tokens,
            )
        )

    async def analyze(
        self,
        *,
        template_image_path: str | None = None,
        template_file_refs: list[dict[str, Any]] | None = None,
    ) -> TemplateAnalysisResult:
        result = TemplateAnalysisResult()
        if template_image_path:
            result = self._merge(
                result, await self.analyze_path(self._resolve_path(template_image_path))
            )
        for item in template_file_refs or []:
            path_value = str(item.get("path") or item.get("file_path") or "").strip()
            if not path_value:
                continue
            result = self._merge(result, await self.analyze_path(self._resolve_path(path_value)))
        return result

    async def analyze_path(self, path: Path) -> TemplateAnalysisResult:
        if not path.exists():
            raise FileNotFoundError(f"Template file not found: {path}")
        suffix = path.suffix.lower()
        if suffix in _IMAGE_EXTENSIONS:
            return await self._analyze_image(path)
        if suffix == ".pptx":
            return await self._analyze_pptx_template(path)
        if suffix == ".pdf":
            return await self._analyze_pdf_template(path)
        return await self._analyze_textual_template(path, "")

    async def _analyze_pdf_template(self, path: Path) -> TemplateAnalysisResult:
        extracted_text = self._extract_pdf_text(path)
        visual_result = TemplateAnalysisResult()
        with tempfile.TemporaryDirectory(prefix="ppt-pdf-preview-") as temp_dir:
            preview_paths = self._render_pdf_preview_images(path, Path(temp_dir))
            if preview_paths:
                visual_result = await self._analyze_preview_images(preview_paths)
        textual_result = await self._analyze_textual_template(path, extracted_text)
        return self._merge(visual_result, textual_result)

    async def _analyze_pptx_template(self, path: Path) -> TemplateAnalysisResult:
        extracted_text = self._extract_pptx_text(path)
        visual_result = TemplateAnalysisResult()
        with tempfile.TemporaryDirectory(prefix="ppt-pptx-preview-") as temp_dir:
            preview_paths = self._render_pptx_preview_images(path, Path(temp_dir))
            if preview_paths:
                visual_result = await self._analyze_preview_images(preview_paths)
        textual_result = await self._analyze_textual_template(path, extracted_text)
        return self._merge(visual_result, textual_result)

    async def _analyze_image(self, path: Path) -> TemplateAnalysisResult:
        image_bytes = path.read_bytes()
        if not image_bytes:
            return TemplateAnalysisResult()
        mime_type = self._guess_image_content_type(path)
        image_b64 = base64.b64encode(image_bytes).decode("utf-8")
        style_system, style_user = PptPromptManager.reference_style_extraction()
        layout_system, layout_user = PptPromptManager.layout_extraction()
        style_data = await self._vision_json_complete(
            style_system, style_user, mime_type, image_b64
        )
        layout_data = await self._vision_json_complete(
            layout_system, layout_user, mime_type, image_b64
        )
        style_prompt = self._normalize_text_field(style_data.get("style_prompt"))
        palette_hint = self._normalize_text_field(style_data.get("palette_hint"))
        composition_hint = self._normalize_text_field(style_data.get("composition_hint"))
        layout_prompt = self._normalize_text_field(layout_data.get("layout_prompt"))
        layout_regions = self._normalize_text_field(layout_data.get("layout_regions"))
        return TemplateAnalysisResult(
            reference_style_prompt="\n".join(
                part
                for part in [
                    style_prompt,
                    f"Palette hint: {palette_hint}" if palette_hint else "",
                    f"Composition hint: {composition_hint}" if composition_hint else "",
                ]
                if part
            ).strip(),
            reference_layout_prompt="\n".join(
                part
                for part in [
                    layout_prompt,
                    f"Layout regions: {layout_regions}" if layout_regions else "",
                ]
                if part
            ).strip(),
        )

    async def _analyze_textual_template(
        self, path: Path, extracted_text: str | None
    ) -> TemplateAnalysisResult:
        trimmed_text = self._trim_text(extracted_text or "", 12000)
        content_system, content_user = PptPromptManager.file_content_extraction(
            file_name=path.name,
            file_text=trimmed_text or "No extracted text available.",
        )
        data = await self._json_complete(content_system, content_user)
        content_prompt = self._normalize_text_field(data.get("content_prompt"))
        key_sections = self._normalize_text_field(data.get("key_sections"))
        content_value = "\n".join(
            part
            for part in [content_prompt, f"Key sections: {key_sections}" if key_sections else ""]
            if part
        ).strip()
        return TemplateAnalysisResult(
            reference_content_prompt=content_value,
            extracted_text=trimmed_text,
        )

    async def _analyze_preview_images(self, preview_paths: list[Path]) -> TemplateAnalysisResult:
        page_results: list[tuple[Path, TemplateAnalysisResult]] = []
        for path in preview_paths[:_PREVIEW_MAX_PAGES]:
            try:
                page_result = await self._analyze_image(path)
                if page_result.reference_style_prompt or page_result.reference_layout_prompt:
                    page_results.append((path, page_result))
            except Exception as exc:
                logger.warning(f"Failed to analyze preview image {path}: {exc}")
        if not page_results:
            return TemplateAnalysisResult()
        if len(page_results) == 1:
            return page_results[0][1]

        merged = TemplateAnalysisResult()
        for _, page_result in page_results:
            merged = self._merge(merged, page_result)

        synthesized = await self._synthesize_preview_results(page_results)
        return TemplateAnalysisResult(
            reference_style_prompt=(
                synthesized.reference_style_prompt or merged.reference_style_prompt
            ),
            reference_layout_prompt=(
                synthesized.reference_layout_prompt or merged.reference_layout_prompt
            ),
        )

    async def _synthesize_preview_results(
        self, page_results: list[tuple[Path, TemplateAnalysisResult]]
    ) -> TemplateAnalysisResult:
        page_findings = [
            {
                "page_index": index,
                "preview_file": path.name,
                "style_prompt": result.reference_style_prompt,
                "layout_prompt": result.reference_layout_prompt,
            }
            for index, (path, result) in enumerate(page_results, start=1)
        ]
        system_prompt, user_prompt = PptPromptManager.visual_synthesis(
            page_findings_json=json.dumps(page_findings, ensure_ascii=False, indent=2)
        )
        try:
            data = await self._json_complete(system_prompt, user_prompt)
        except Exception as exc:
            logger.warning(f"Failed to synthesize preview findings: {exc}")
            return TemplateAnalysisResult()
        return TemplateAnalysisResult(
            reference_style_prompt=(data.get("reference_style_prompt") or "").strip(),
            reference_layout_prompt=(data.get("reference_layout_prompt") or "").strip(),
        )

    async def _vision_json_complete(
        self,
        system_prompt: str,
        user_prompt: str,
        mime_type: str,
        image_b64: str,
    ) -> dict[str, Any]:
        kwargs = get_token_limit_kwargs(self.config.model, self.config.max_tokens)
        vision_func = self._vision_client.get_vision_model_func()
        use_json_mode = supports_response_format_json_object(
            self.config.model,
            self.config.binding,
            self.config.base_url,
        )
        try:
            raw = await vision_func(
                prompt="",
                messages=self._build_vision_messages(
                    system_prompt,
                    user_prompt if use_json_mode else self._force_json_output(user_prompt),
                    mime_type,
                    image_b64,
                ),
                **({"response_format": {"type": "json_object"}} if use_json_mode else {}),
                temperature=self.config.temperature,
                **kwargs,
            )
        except Exception as exc:
            if not use_json_mode or not self._should_retry_without_json_mode(exc):
                raise
            logger.warning(
                f"Vision model {self.config.model} does not support "
                "response_format=json_object; retrying without JSON mode"
            )
            raw = await vision_func(
                prompt="",
                messages=self._build_vision_messages(
                    system_prompt,
                    self._force_json_output(user_prompt),
                    mime_type,
                    image_b64,
                ),
                temperature=self.config.temperature,
                **kwargs,
            )
        data = self._extract_json(raw)
        if not data:
            raise ValueError("Failed to parse JSON output from template-analysis vision model")
        return data

    async def _json_complete(self, system_prompt: str, user_prompt: str) -> dict[str, Any]:
        kwargs = get_token_limit_kwargs(self.config.model, self.config.max_tokens)
        use_json_mode = supports_response_format_json_object(
            self.config.model,
            self.config.binding,
            self.config.base_url,
        )
        try:
            raw = await llm_complete(
                prompt=user_prompt if use_json_mode else self._force_json_output(user_prompt),
                system_prompt=system_prompt,
                model=self.config.model,
                api_key=self.config.api_key,
                base_url=self.config.base_url,
                binding=self.config.binding,
                temperature=self.config.temperature,
                **({"response_format": {"type": "json_object"}} if use_json_mode else {}),
                **kwargs,
            )
        except Exception as exc:
            if not use_json_mode or not self._should_retry_without_json_mode(exc):
                raise
            logger.warning(
                f"Model {self.config.model} does not support "
                "response_format=json_object; retrying without JSON mode"
            )
            raw = await llm_complete(
                prompt=self._force_json_output(user_prompt),
                system_prompt=system_prompt,
                model=self.config.model,
                api_key=self.config.api_key,
                base_url=self.config.base_url,
                binding=self.config.binding,
                temperature=self.config.temperature,
                **kwargs,
            )
        data = self._extract_json(raw)
        if not data:
            raise ValueError("Failed to parse JSON output from template-analysis model")
        return data

    def _build_vision_messages(
        self,
        system_prompt: str,
        user_prompt: str,
        mime_type: str,
        image_b64: str,
    ) -> list[dict[str, Any]]:
        return [
            {"role": "system", "content": system_prompt},
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": user_prompt},
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:{mime_type};base64,{image_b64}"},
                    },
                ],
            },
        ]

    def _force_json_output(self, user_prompt: str) -> str:
        return (
            f"{user_prompt.rstrip()}\n\n"
            "IMPORTANT: Return a single valid JSON object only. "
            "Do not use markdown fences. Do not add explanations before or after the JSON."
        )

    def _normalize_text_field(self, value: Any) -> str:
        if value is None:
            return ""
        if isinstance(value, str):
            return value.strip()
        if isinstance(value, dict):
            parts = []
            for key, nested in value.items():
                nested_text = self._normalize_text_field(nested)
                if nested_text:
                    parts.append(f"{key}: {nested_text}")
            return "; ".join(parts).strip()
        if isinstance(value, (list, tuple, set)):
            parts = [self._normalize_text_field(item) for item in value]
            return "; ".join(part for part in parts if part).strip()
        return str(value).strip()

    def _should_retry_without_json_mode(self, exc: Exception) -> bool:
        message = str(exc).lower()
        if "json_object" not in message:
            return False
        return any(
            marker in message
            for marker in ("not supported", "not valid", "invalidparameter", "unsupported")
        )

    def _resolve_path(self, value: str) -> Path:
        normalized = (value or "").strip()
        if normalized.startswith("/api/outputs/"):
            normalized = normalized.replace("/api/outputs/", "", 1)
        if normalized.startswith("/"):
            return Path(normalized)
        return self.project_root / "data" / "user" / normalized

    def _merge(
        self, base: TemplateAnalysisResult, incoming: TemplateAnalysisResult
    ) -> TemplateAnalysisResult:
        return TemplateAnalysisResult(
            reference_style_prompt=self._join_unique(
                base.reference_style_prompt, incoming.reference_style_prompt
            ),
            reference_layout_prompt=self._join_unique(
                base.reference_layout_prompt, incoming.reference_layout_prompt
            ),
            reference_content_prompt=self._join_unique(
                base.reference_content_prompt, incoming.reference_content_prompt
            ),
            extracted_text=self._join_unique(base.extracted_text, incoming.extracted_text),
        )

    def _join_unique(self, *parts: str) -> str:
        values: list[str] = []
        seen: set[str] = set()
        for part in parts:
            text = (part or "").strip()
            if not text or text in seen:
                continue
            seen.add(text)
            values.append(text)
        return "\n\n".join(values).strip()

    def _render_pdf_preview_images(self, path: Path, out_dir: Path) -> list[Path]:
        try:
            import pypdfium2 as pdfium
        except Exception as exc:
            logger.warning(f"pypdfium2 unavailable; cannot render PDF previews: {exc}")
            return []

        out_dir.mkdir(parents=True, exist_ok=True)
        preview_paths: list[Path] = []
        try:
            document = pdfium.PdfDocument(str(path))
        except Exception as exc:
            logger.warning(f"Failed to open PDF for preview rendering: {exc}")
            return []

        page_count = min(len(document), _PREVIEW_MAX_PAGES)
        for index in range(page_count):
            png_path = out_dir / f"{path.stem}_page_{index + 1}.png"
            page = None
            try:
                page = document[index]
                pil_image = page.render(scale=2.0).to_pil()
                pil_image.save(png_path, format="PNG")
                preview_paths.append(png_path)
            except Exception as exc:
                logger.warning(f"Failed to render PDF preview for {path} page {index + 1}: {exc}")
            finally:
                if page is not None:
                    try:
                        page.close()
                    except Exception:
                        pass
        try:
            document.close()
        except Exception:
            pass
        return preview_paths

    def _render_pptx_preview_images(self, path: Path, out_dir: Path) -> list[Path]:
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except Exception as exc:
            logger.warning(f"python-pptx unavailable for PPTX preview rendering: {exc}")
            return []

        out_dir.mkdir(parents=True, exist_ok=True)
        try:
            presentation = Presentation(str(path))
        except Exception as exc:
            logger.warning(f"Failed to open PPTX for preview rendering: {exc}")
            return []

        slide_width = int(presentation.slide_width or 0)
        slide_height = int(presentation.slide_height or 0)
        if not slide_width or not slide_height:
            return []
        scale = _PREVIEW_WIDTH / slide_width
        canvas_height = max(1, int(slide_height * scale))
        preview_paths: list[Path] = []

        for index, slide in enumerate(presentation.slides, start=1):
            if index > _PREVIEW_MAX_PAGES:
                break
            image = Image.new(
                "RGB",
                (_PREVIEW_WIDTH, canvas_height),
                self._extract_background_color(slide) or _DEFAULT_BG,
            )
            draw = ImageDraw.Draw(image)
            for shape in slide.shapes:
                self._draw_pptx_shape(
                    draw=draw,
                    canvas=image,
                    shape=shape,
                    scale=scale,
                    shape_type_enum=MSO_SHAPE_TYPE,
                )
            out_path = out_dir / f"{path.stem}_slide_{index}.png"
            image.save(out_path, format="PNG")
            preview_paths.append(out_path)
        return preview_paths

    def _draw_pptx_shape(
        self,
        *,
        draw: ImageDraw.ImageDraw,
        canvas: Image.Image,
        shape: Any,
        scale: float,
        shape_type_enum: Any,
    ) -> None:
        left = int(getattr(shape, "left", 0) * scale)
        top = int(getattr(shape, "top", 0) * scale)
        width = max(1, int(getattr(shape, "width", 1) * scale))
        height = max(1, int(getattr(shape, "height", 1) * scale))
        bounds = (left, top, left + width, top + height)

        if getattr(shape, "shape_type", None) == getattr(shape_type_enum, "GROUP", None):
            for child in getattr(shape, "shapes", []) or []:
                self._draw_pptx_shape(
                    draw=draw,
                    canvas=canvas,
                    shape=child,
                    scale=scale,
                    shape_type_enum=shape_type_enum,
                )
            return

        if getattr(shape, "shape_type", None) == getattr(shape_type_enum, "PICTURE", None):
            if self._paste_picture(canvas, shape, bounds):
                draw.rectangle(bounds, outline=_DEFAULT_SHAPE_LINE, width=2)
            else:
                draw.rectangle(bounds, fill=(226, 232, 240), outline=_DEFAULT_SHAPE_LINE, width=2)
            return

        fill_color = self._extract_fill_color(shape) or _DEFAULT_SHAPE_FILL
        line_color = self._extract_line_color(shape) or _DEFAULT_SHAPE_LINE
        draw.rectangle(bounds, fill=fill_color, outline=line_color, width=2)

        if getattr(shape, "has_chart", False):
            self._draw_chart_placeholder(draw, bounds)
            return
        if getattr(shape, "has_table", False):
            self._draw_table_placeholder(draw, bounds)
            return
        if getattr(shape, "has_text_frame", False) or str(getattr(shape, "text", "") or "").strip():
            self._draw_text(draw, bounds, str(getattr(shape, "text", "") or "").strip())

    def _paste_picture(
        self, canvas: Image.Image, shape: Any, bounds: tuple[int, int, int, int]
    ) -> bool:
        try:
            blob = shape.image.blob
            picture = Image.open(BytesIO(blob)).convert("RGB")
            resized = picture.resize((bounds[2] - bounds[0], bounds[3] - bounds[1]))
            canvas.paste(resized, (bounds[0], bounds[1]))
            return True
        except Exception:
            return False

    def _draw_chart_placeholder(
        self, draw: ImageDraw.ImageDraw, bounds: tuple[int, int, int, int]
    ) -> None:
        left, top, right, bottom = bounds
        width = right - left
        height = bottom - top
        bar_width = max(10, width // 8)
        baseline = bottom - 16
        for index, ratio in enumerate((0.35, 0.55, 0.8, 0.6)):
            x0 = left + 18 + index * (bar_width + 12)
            x1 = min(right - 18, x0 + bar_width)
            y1 = baseline
            y0 = max(top + 18, int(baseline - height * ratio))
            draw.rectangle((x0, y0, x1, y1), fill=(59, 130, 246), outline=None)
        draw.line((left + 12, baseline, right - 12, baseline), fill=_DEFAULT_TEXT, width=2)

    def _draw_table_placeholder(
        self, draw: ImageDraw.ImageDraw, bounds: tuple[int, int, int, int]
    ) -> None:
        left, top, right, bottom = bounds
        rows = 4
        cols = 3
        for row in range(1, rows):
            y = top + ((bottom - top) * row) // rows
            draw.line((left, y, right, y), fill=_DEFAULT_SHAPE_LINE, width=1)
        for col in range(1, cols):
            x = left + ((right - left) * col) // cols
            draw.line((x, top, x, bottom), fill=_DEFAULT_SHAPE_LINE, width=1)

    def _draw_text(
        self,
        draw: ImageDraw.ImageDraw,
        bounds: tuple[int, int, int, int],
        text: str,
    ) -> None:
        if not text:
            return
        left, top, right, bottom = bounds
        font = ImageFont.load_default()
        margin = 10
        max_chars = max(10, (right - left - margin * 2) // 7)
        lines: list[str] = []
        for raw_line in text.splitlines():
            stripped = raw_line.strip()
            if not stripped:
                continue
            while len(stripped) > max_chars:
                lines.append(stripped[:max_chars])
                stripped = stripped[max_chars:]
            lines.append(stripped)
        y = top + margin
        for line in lines[:8]:
            if y > bottom - 18:
                break
            draw.text((left + margin, y), line, fill=_DEFAULT_TEXT, font=font)
            y += 18

    def _extract_background_color(self, slide: Any) -> tuple[int, int, int] | None:
        try:
            return self._extract_rgb(getattr(slide.background.fill.fore_color, "rgb", None))
        except Exception:
            return None

    def _extract_fill_color(self, shape: Any) -> tuple[int, int, int] | None:
        try:
            fill = shape.fill
            if getattr(fill, "type", None) is None:
                return None
            return self._extract_rgb(getattr(fill.fore_color, "rgb", None))
        except Exception:
            return None

    def _extract_line_color(self, shape: Any) -> tuple[int, int, int] | None:
        try:
            line = shape.line
            return self._extract_rgb(getattr(line.color, "rgb", None))
        except Exception:
            return None

    def _extract_rgb(self, value: Any) -> tuple[int, int, int] | None:
        if value is None:
            return None
        text = str(value).strip().replace("RGBColor(", "").replace(")", "").replace(",", "")
        text = text.lstrip("#")
        if len(text) != 6:
            return None
        try:
            return tuple(int(text[index : index + 2], 16) for index in (0, 2, 4))
        except Exception:
            return None

    def _extract_pptx_text(self, path: Path) -> str:
        try:
            from pptx import Presentation
        except Exception as exc:
            logger.warning(f"python-pptx unavailable for template analysis: {exc}")
            return ""

        try:
            presentation = Presentation(str(path))
        except Exception as exc:
            logger.warning(f"Failed to read PPTX template text: {exc}")
            return ""

        blocks: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                text = str(getattr(shape, "text", "") or "").strip()
                if text:
                    texts.append(text)
            if texts:
                blocks.append(f"## Slide {index}\n" + "\n".join(texts[:12]))
        return "\n\n".join(blocks)

    def _extract_pdf_text(self, path: Path) -> str:
        for module_name, class_name in (("pypdf", "PdfReader"), ("PyPDF2", "PdfReader")):
            try:
                module = __import__(module_name, fromlist=[class_name])
                reader_cls = getattr(module, class_name)
                reader = reader_cls(str(path))
                pages: list[str] = []
                for page in reader.pages[:8]:
                    text = (page.extract_text() or "").strip()
                    if text:
                        pages.append(text)
                return "\n\n".join(pages)
            except Exception:
                continue
        logger.warning("No PDF reader available for template analysis")
        return ""

    def _guess_image_content_type(self, image_path: Path) -> str:
        suffix = image_path.suffix.lower()
        if suffix == ".png":
            return "image/png"
        if suffix == ".webp":
            return "image/webp"
        return "image/jpeg"

    def _extract_json(self, text: str) -> dict[str, Any]:
        if not text:
            return {}
        try:
            data = json.loads(text)
            return data if isinstance(data, dict) else {}
        except json.JSONDecodeError:
            pass
        match = re.search(r"```json\s*([\s\S]*?)\s*```", text, flags=re.IGNORECASE)
        if match:
            try:
                data = json.loads(match.group(1).strip())
                return data if isinstance(data, dict) else {}
            except Exception:
                return {}
        start = text.find("{")
        end = text.rfind("}")
        if start != -1 and end != -1 and end > start:
            try:
                data = json.loads(text[start : end + 1])
                return data if isinstance(data, dict) else {}
            except Exception:
                return {}
        return {}

    def _trim_text(self, value: str, max_chars: int) -> str:
        text = (value or "").strip()
        if len(text) <= max_chars:
            return text
        return text[: max_chars - 1].rstrip() + "…"


__all__ = ["TemplateAnalysisResult", "TemplateAnalyzer"]
